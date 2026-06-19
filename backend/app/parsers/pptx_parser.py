from pathlib import Path

from pptx import Presentation

from app.parsers.base import ParsedBlock, ParsedDocument
from app.parsers.pdf_parser import _guess_category


def parse_pptx(path: Path) -> ParsedDocument:
    prs = Presentation(path)
    blocks: list[ParsedBlock] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        text = "\n".join(parts)
        if notes:
            text = f"{text}\n\n[备注] {notes}" if text else f"[备注] {notes}"
        if text.strip():
            blocks.append(
                ParsedBlock(
                    text=text.strip(),
                    page=i,
                    section=f"幻灯片 {i}",
                    metadata={"type": "slide"},
                )
            )
    return ParsedDocument(blocks=blocks, doc_type="pptx", category=_guess_category(path.name))
