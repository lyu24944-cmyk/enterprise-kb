"""生成演示文档并入库到知识库。"""
from __future__ import annotations

import asyncio
import shutil
import sys
import uuid
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
BACKEND = ROOT / "backend"
sys.path.insert(0, str(BACKEND))

DEMO_DIR = ROOT / "data" / "demo"


def create_handbook_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("星辰科技有限公司员工手册（2024版）", 0)
    doc.add_heading("第一章 总则", 1)
    doc.add_paragraph(
        "本手册适用于星辰科技有限公司全体正式员工。员工入职即视为同意遵守本手册及相关规章制度。"
    )
    doc.add_heading("第二章 工作时间与考勤", 1)
    doc.add_paragraph("工作时间：周一至周五 9:00-18:00，午休 12:00-13:00。")
    doc.add_paragraph("迟到 30 分钟以内累计三次记一次旷工；旷工一天扣除当日工资并通报批评。")
    doc.add_heading("第三章 休假制度", 1)
    doc.add_paragraph("3.1 法定节假日：按国家规定执行。")
    doc.add_paragraph(
        "3.2 带薪年假：员工累计工作满 1 年不满 10 年的，年假 5 天；"
        "满 10 年不满 20 年的，年假 10 天；满 20 年的，年假 15 天。"
        "新入职员工当年度年假按在职月份折算，不足 1 个月不享受。"
    )
    doc.add_paragraph("3.3 事假：需提前 1 个工作日通过 OA 提交申请，经直属主管审批。")
    doc.add_paragraph("3.4 病假：需提供二级以上医院证明，3 天以上需 HR 备案。")
    doc.add_paragraph("3.5 婚假：依法享受 3 天，晚婚额外奖励假 7 天（以公司政策为准）。")
    doc.add_heading("第四章 薪酬福利", 1)
    doc.add_paragraph("工资于每月 10 日发放。年终奖根据公司效益及个人绩效确定，不承诺固定金额。")
    doc.add_heading("第五章 保密义务", 1)
    doc.add_paragraph("员工对在职期间知悉的商业秘密、技术资料负有保密义务，离职后两年内仍有效。")
    doc.save(path)


def create_contract_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("软件开发服务合同", 0)
    doc.add_paragraph("甲方（委托方）：星辰科技有限公司")
    doc.add_paragraph("乙方（服务方）：____________")
    doc.add_heading("第一条 服务内容", 1)
    doc.add_paragraph("乙方向甲方提供企业知识库系统定制开发服务，交付源代码及部署文档。")
    doc.add_heading("第二条 合同金额与付款", 1)
    doc.add_paragraph(
        "合同总金额人民币 500,000 元。甲方于合同签订后 5 日内支付 30% 预付款；"
        "验收合格后 30 日内支付 60%；质保期满一年后支付尾款 10%。"
    )
    doc.add_paragraph("若甲方逾期付款超过 15 日，乙方有权暂停服务且不承担违约责任。")
    doc.add_heading("第三条 知识产权", 1)
    doc.add_paragraph(
        "项目交付的全部成果（含源代码、文档、模型权重衍生成果）知识产权归甲方独占所有。"
        "乙方不得将本项目成果用于其他客户或开源发布。"
    )
    doc.add_heading("第四条 违约责任", 1)
    doc.add_paragraph(
        "乙方延期交付每日按合同总额 0.5% 支付违约金，无上限。"
        "甲方延期付款每日按欠款额 0.1% 支付违约金，累计不超过欠款额的 10%。"
    )
    doc.add_heading("第五条 保密", 1)
    doc.add_paragraph("双方对合同内容及履行中知悉的信息保密，期限为合同终止后五年。")
    doc.add_heading("第六条 争议解决", 1)
    doc.add_paragraph("争议提交甲方所在地人民法院诉讼解决。诉讼费用由败诉方承担。")
    doc.add_heading("第七条 合同终止", 1)
    doc.add_paragraph(
        "甲方有权提前 30 日书面通知单方解除合同，乙方已发生成本不予退还，"
        "且需向甲方移交全部阶段性成果。"
    )
    doc.save(path)


def create_training_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "星辰科技 · 新员工制度培训"
    slide.placeholders[1].text = "HR 部 · 2024"

    for title, body in [
        ("考勤要点", "9:00 前打卡\n迟到需提交说明\n外勤需提前报备"),
        ("请假流程", "1. OA 提交申请\n2. 主管审批\n3. HR 备案\n4. 销假确认"),
        ("年假规则", "满1年：5天\n满10年：10天\n满20年：15天\n入职当年按月份折算"),
        ("保密要求", "禁止外传客户资料\n离职交接需签署保密承诺书"),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.text = body.split("\n")[0]
        for line in body.split("\n")[1:]:
            p = tf.add_paragraph()
            p.text = line

    prs.save(path)


def create_leave_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "请假类型"
    ws.append(["类型", "提前申请", "审批人", "证明材料", "备注"])
    ws.append(["年假", "1个工作日", "直属主管", "无", "按手册第三章执行"])
    ws.append(["事假", "1个工作日", "直属主管", "无", "不带薪"])
    ws.append(["病假", "可事后补", "直属主管+HR", "医院证明", "3天以上需HR备案"])
    ws.append(["婚假", "5个工作日", "HR", "结婚证", "依法3天+公司奖励"])
    ws.append(["产假", "按法规", "HR", "相关证明", "按国家和地方政策"])

    ws2 = wb.create_sheet("审批时限")
    ws2.append(["环节", "时限", "说明"])
    ws2.append(["主管审批", "2个工作日", "超时自动提醒"])
    ws2.append(["HR备案", "1个工作日", "病假、产假必须"])
    wb.save(path)


def create_handbook_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    text = (
        "星辰科技有限公司员工手册（2024版）\n\n第三章 休假制度\n\n"
        "3.2 带薪年假：员工累计工作满 1 年不满 10 年的，年假 5 天；"
        "满 10 年不满 20 年的，年假 10 天；满 20 年的，年假 15 天。"
    )
    page = doc.new_page()
    page.insert_text((50, 50), text, fontsize=11)
    doc.save(path)
    doc.close()


def generate_all() -> list[Path]:
    DEMO_DIR.mkdir(parents=True, exist_ok=True)
    builders = {
        "员工手册2024.docx": create_handbook_docx,
        "员工手册2024.pdf": create_handbook_pdf,
        "软件开发服务合同.docx": create_contract_docx,
        "公司制度培训.pptx": create_training_pptx,
        "请假制度表.xlsx": create_leave_xlsx,
    }
    paths: list[Path] = []
    for name, fn in builders.items():
        p = DEMO_DIR / name
        fn(p)
        paths.append(p)
        print(f"  生成: {p.name}")
    return paths


async def ingest_all(paths: list[Path]) -> None:
    from sqlalchemy import select

    from app.core.config import get_settings
    from app.models.database import Document
    from app.services.db_service import SessionLocal, ingest_document, init_db

    await init_db()
    settings = get_settings()
    settings.upload_path.mkdir(parents=True, exist_ok=True)

    async with SessionLocal() as session:
        for p in paths:
            result = await session.execute(select(Document).where(Document.filename == p.name))
            if result.scalar_one_or_none():
                print(f"  跳过已存在: {p.name}")
                continue
            target = settings.upload_path / f"{uuid.uuid4().hex}_{p.name}"
            shutil.copy(p, target)
            doc = await ingest_document(session, target, p.name)
            print(f"  入库: {doc.filename} ({doc.chunk_count} chunks)")


def main() -> None:
    print("=== 生成演示文档 ===")
    paths = generate_all()
    print("\n=== 入库到向量库（首次会下载 Embedding 模型）===")
    asyncio.run(ingest_all(paths))
    print("\n完成！启动后端和前端后即可提问：公司年假多少天？")


if __name__ == "__main__":
    main()
